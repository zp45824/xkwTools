from pptx import Presentation
import os
import os.path
import win32com.client

PAGE = -1

def del_ppt_max_page(file):
    # 读取ppt
    prs = Presentation(file)
    # 查看一共几页
    slides = prs.slides
    print(slides)
    number_pages = len(slides)
    print('file={},处理前总页数={}'.format(file, number_pages))
    # 删除最后一页
    rId = prs.slides._sldIdLst[PAGE].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[PAGE]

    # rId = prs.slides._sldIdLst[PAGE].rId
    # prs.part.drop_rel(rId)
    # del prs.slides._sldIdLst[PAGE]
    print('file={},处理后总页数={}'.format(file, len(slides)))

    # 保存新的ppt
    prs.save(file)


def list_dir_more(CurPath=os.getcwd(), file_list=[]):
    FileList = os.listdir(CurPath)
    # print FileList
    for File in FileList:
        SubPath = CurPath + '\\' + File
        if os.path.isdir(SubPath):
            list_dir_more(SubPath, file_list)
        else:
            file_list.append(SubPath)

    # print('file_list={}'.format(file_list))
    return file_list


if __name__ == "__main__":
    powerpoint = win32com.client.Dispatch('PowerPoint.Application')
    win32com.client.gencache.EnsureDispatch('PowerPoint.Application')

    powerpoint.Visible = True
    dir = input("请输入路径")
    file_list = list_dir_more(dir)
    count = 1
    for file in file_list:
        if file[-4:] == "pptx":
            del_ppt_max_page(file)
            print('正在处理第 {}/{} 个文件：{}'.format(count, len(file_list), file))
        if file[-3:] == "ppt":
            ppt1 = powerpoint.Presentations.Open(file)
            ppt1.SaveAs(file[:-4] + '.pptx')
            ppt1.Close()
            os.remove(file)
            file = file + 'x'
            del_ppt_max_page(file)
            print('正在处理第 {}/{} 个文件：{}'.format(count, len(file_list), file))

        count += 1




